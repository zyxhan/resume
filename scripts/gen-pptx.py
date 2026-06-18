from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ======== Color Palette ========
INK = RGBColor(0x0a, 0x0a, 0x0b)
INK_TINT = RGBColor(0x18, 0x18, 0x1a)
PAPER = RGBColor(0xf1, 0xef, 0xea)
PAPER_TINT = RGBColor(0xe8, 0xe5, 0xde)
MUTED_LIGHT = RGBColor(0x99, 0x96, 0x8e)
MUTED_DARK = RGBColor(0x66, 0x64, 0x60)
WHITE_SOFT = RGBColor(0xfa, 0xf8, 0xf5)
BODY_DARK = RGBColor(0xcc, 0xca, 0xc5)
BODY_LIGHT = RGBColor(0x66, 0x63, 0x5e)
CALLOUT_BG = RGBColor(0x1a, 0x1a, 0x1c)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_tb(slide, left, top, width, height, text="", font_size=Pt(14),
           color=INK, bold=False, font_name='Microsoft YaHei',
           alignment=PP_ALIGN.LEFT, line_sp=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_sp != 1.0:
        p.line_spacing = Pt(font_size.pt * line_sp)
    return tf


def add_p(tf, text, font_size=Pt(14), color=INK, bold=False,
          font_name='Microsoft YaHei', alignment=PP_ALIGN.LEFT,
          sp_before=Pt(4), sp_after=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = sp_before
    p.space_after = sp_after
    return p


def kicker(slide, text, left, top, width, dark=False):
    c = MUTED_LIGHT if dark else MUTED_DARK
    return add_tb(slide, left, top, width, Inches(0.35), text,
                  Pt(10), c, False, 'Consolas')


def h1(slide, text, left, top, width, dark=False):
    c = PAPER if dark else INK
    return add_tb(slide, left, top, width, Inches(1.1), text,
                  Pt(42), c, True, 'Microsoft YaHei', line_sp=1.1)


def h2(slide, text, left, top, width, dark=False):
    c = MUTED_LIGHT if dark else MUTED_DARK
    return add_tb(slide, left, top, width, Inches(0.7), text,
                  Pt(24), c, False, 'Microsoft YaHei', line_sp=1.2)


def lead(slide, text, left, top, width, dark=False):
    c = PAPER if dark else INK
    return add_tb(slide, left, top, width, Inches(0.9), text,
                  Pt(18), c, False, 'Microsoft YaHei', line_sp=1.5)


def body(slide, text, left, top, width, height=Inches(0.6), dark=False):
    c = BODY_DARK if dark else BODY_LIGHT
    return add_tb(slide, left, top, width, height, text,
                  Pt(13), c, False, 'Microsoft YaHei', line_sp=1.6)


def line(slide, left, top, width, dark=False):
    c = MUTED_LIGHT if dark else MUTED_DARK
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()


def mono(slide, text, left, top, width, dark=False, align=PP_ALIGN.LEFT):
    c = MUTED_LIGHT if dark else MUTED_DARK
    return add_tb(slide, left, top, width, Inches(0.3), text,
                  Pt(9), c, False, 'Consolas', align)


def card(slide, left, top, width, label, val, note, dark=False):
    lc = MUTED_LIGHT if dark else MUTED_DARK
    vc = PAPER if dark else INK
    nc = BODY_DARK if dark else MUTED_DARK
    tf = add_tb(slide, left, top, width, Inches(0.25), label,
                Pt(9), lc, False, 'Consolas')
    add_p(tf, val, Pt(28), vc, True, sp_before=Pt(6))
    add_p(tf, note, Pt(13), nc, False, sp_before=Pt(8))
    return tf


def chrome(slide, left_text, right_text=None, dark=False, page=""):
    mono(slide, left_text, Inches(1), Inches(0.6), Inches(9), dark)
    if page:
        mono(slide, page, Inches(10.5), Inches(0.6), Inches(2), dark, PP_ALIGN.RIGHT)
    elif right_text:
        mono(slide, right_text, Inches(10.5), Inches(0.6), Inches(2), dark, PP_ALIGN.RIGHT)


def foot(slide, left_text, right_text="— · —", dark=False):
    mono(slide, left_text, Inches(1), Inches(6.8), Inches(6), dark)
    mono(slide, right_text, Inches(8), Inches(6.8), Inches(4.5), dark, PP_ALIGN.RIGHT)


def callout_box(slide, left, top, width, height, quote, src, dark=False):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CALLOUT_BG if dark else WHITE_SOFT
    bg.line.fill.background()
    qc = PAPER if dark else INK
    sc = MUTED_LIGHT if dark else MUTED_DARK
    add_tb(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.5),
           Inches(0.5), quote, Pt(14), qc, False, 'Microsoft YaHei', line_sp=1.4)
    add_tb(slide, left + Inches(0.25), top + height - Inches(0.4),
           width - Inches(0.5), Inches(0.25), src, Pt(8), sc, False, 'Consolas')


# ============================================================
# SLIDE 1: Cover (hero dark)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, INK)
chrome(sl, "计算机学院党员工作站 · 2025", "述职报告", True)

kicker(sl, "职务调整申请", Inches(1), Inches(1.8), Inches(11), True)
add_tb(sl, Inches(1), Inches(2.4), Inches(11), Inches(1.3), "安全服务部",
       Pt(52), PAPER, True, 'Microsoft YaHei', line_sp=1.0)
h2(sl, "留部述职报告", Inches(1), Inches(3.5), Inches(11), True)
lead(sl, "一年以来，每周查寝、每次挂牌，\n用最日常的工作，打磨最扎实的责任感。",
     Inches(1), Inches(4.5), Inches(8), True)
mono(sl, "述职人：朱轶涵 · 25级软件工程二班 · 山东泰安",
     Inches(1), Inches(5.8), Inches(11), True)
foot(sl, "安全服务部 · 述职报告", "— 2025 —", True)

# ============================================================
# SLIDE 2: Contents (light)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, PAPER)
chrome(sl, "Contents · 目录", page="02 / 14")

kicker(sl, "Overview · 概览", Inches(1), Inches(1.5), Inches(11))
h1(sl, "四件事，四个维度。", Inches(1), Inches(2.0), Inches(11))

items = [
    ("I", "我是谁 —— 一个软件工程学生的自述", "About"),
    ("II", "安全服务部 —— 守护一栋楼的秩序与温度", "Department"),
    ("III", "走过的路 —— 在重复的工作里长出判断力", "Growth"),
    ("IV", "下一步 —— 留下来，做得更好", "Future"),
]
for i, (k, v, m) in enumerate(items):
    y = Inches(3.3 + i * 0.75)
    line(sl, Inches(1), y, Inches(11.3))
    add_tb(sl, Inches(1), y + Pt(6), Inches(0.5), Inches(0.4), k,
           Pt(20), INK, True, 'Microsoft YaHei')
    add_tb(sl, Inches(1.8), y + Pt(6), Inches(7), Inches(0.4), v,
           Pt(15), BODY_LIGHT, False, 'Microsoft YaHei')
    mono(sl, m, Inches(10), y + Pt(6), Inches(2.3), align=PP_ALIGN.RIGHT)

foot(sl, "目录 · 述职结构")

# ============================================================
# SLIDE 3: Chapter Divider - About (hero light)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, PAPER_TINT)
chrome(sl, "第一幕 · 自我介绍", page="03 / 14")

kicker(sl, "Chapter I", Inches(1), Inches(2.2), Inches(5))
add_tb(sl, Inches(1), Inches(2.9), Inches(10), Inches(1.1), "我是谁",
       Pt(60), INK, True, 'Microsoft YaHei', line_sp=1.0)
lead(sl, "一个来自山东泰安的普通人，\n在代码与查寝之间寻找平衡。",
     Inches(1), Inches(4.3), Inches(7))
foot(sl, "第一幕 · 个人画像")

# ============================================================
# SLIDE 4: Personal Intro (dark)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, INK)
chrome(sl, "Profile · 个人画像", page="04 / 14", dark=True)

kicker(sl, "关于我", Inches(1), Inches(1.5), Inches(5.5), True)
add_tb(sl, Inches(1), Inches(2.0), Inches(5.5), Inches(0.7), "朱轶涵",
       Pt(34), PAPER, True, 'Microsoft YaHei')
lead(sl, "25级软件工程二班，山东泰安。\n喜欢听歌、散步——在音乐里找到节奏，\n在走路时想清楚事情。",
     Inches(1), Inches(2.8), Inches(5.5), True)

callout_box(sl, Inches(1), Inches(4.8), Inches(5.5), Inches(1.1),
            "「不张扬，但经得住日复一日的考验。」",
            "— 关于我自己的总结", True)

body(sl, "进入大学这一年，学会的最重要的事不是写代码，而是在重复的日常里保持靠谱——每周的查寝不迟到、不敷衍，每次的任务不拖延、不推诿。",
     Inches(7.2), Inches(1.8), Inches(5), Inches(1.0), True)
body(sl, "我不是那种站在台上光芒万丈的人。但交给我的事，我会一件一件做完，一件一件做好。软件工程教会我逻辑，安全服务部教会我责任。",
     Inches(7.2), Inches(3.2), Inches(5), Inches(1.0), True)

tags = [("班级", "25软工二班"), ("家乡", "山东·泰安"), ("兴趣", "听歌·散步")]
for i, (lbl, val) in enumerate(tags):
    x = Inches(7.2 + i * 1.8)
    mono(sl, lbl, x, Inches(4.8), Inches(1.6), True)
    add_tb(sl, x, Inches(5.1), Inches(1.6), Inches(0.3), val,
           Pt(15), PAPER, True, 'Microsoft YaHei')

foot(sl, "Page 04 · 个人画像", dark=True)

# ============================================================
# SLIDE 5: Work Attitude (light)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, PAPER)
chrome(sl, "Attitude · 工作态度", page="05 / 14")

kicker(sl, "How I Work · 我的方式", Inches(1), Inches(1.5), Inches(11))
h1(sl, "三句话，概括我的工作底色。", Inches(1), Inches(2.0), Inches(11))

cards = [
    ("01 · Consistency", "准时到场", "每周查寝提前十分钟到，从不缺席。一学期下来，不需要催也不需要提醒，这是我的底线。"),
    ("02 · Ownership", "闭环做事", "领任务从\"收到\"开始，到\"已完成\"结束。挂牌从人员分配到照片收取，每一步都盯到底。"),
    ("03 · Collaboration", "协作沟通", "查寝不是一个人的事。跟安全部对接问题表、跟宣传部确认优差寝照片，都需要主动沟通。"),
]
for i, (lbl, val, note) in enumerate(cards):
    x = Inches(1 + i * 3.9)
    card(sl, x, Inches(3.0), Inches(3.5), lbl, val, note)

foot(sl, "Page 05 · 工作态度")

# ============================================================
# SLIDE 6: Chapter - Department (hero dark)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, INK)
chrome(sl, "第二幕 · 部门认知", page="06 / 14", dark=True)

kicker(sl, "Chapter II", Inches(1), Inches(2.2), Inches(5), True)
add_tb(sl, Inches(1), Inches(2.9), Inches(10), Inches(1.1), "安全服务部",
       Pt(56), PAPER, True, 'Microsoft YaHei', line_sp=1.0)
lead(sl, "把一栋宿舍楼守好，让每一次查寝都有据可循。",
     Inches(1), Inches(4.3), Inches(7), True)
foot(sl, "第二幕 · 安全服务部", dark=True)

# ============================================================
# SLIDE 7: Three Core Duties (light)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, PAPER)
chrome(sl, "Core Duties · 核心职责", page="07 / 14")

kicker(sl, "What We Do · 三个关键词", Inches(1), Inches(1.5), Inches(11))
h1(sl, "安全服务部，三件大事。", Inches(1), Inches(2.0), Inches(11))

duties = [
    ("例行查寝", "每周一次，覆盖宿舍卫生与安全。提前通知、准时到场、规范打分、及时反馈。"),
    ("党员查寝", "每月一次，组织党员开展主题查寝。负责签到、分组、心得收取与归档。"),
    ("党员挂牌", "每学期一次，协调人员分配、发放牌子与床铺贴、维持现场纪律、收取照片归档。"),
]
for i, (title, desc) in enumerate(duties):
    x = Inches(1 + i * 3.9)
    cbg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.0), Inches(3.5), Inches(2.6))
    cbg.fill.solid()
    cbg.fill.fore_color.rgb = WHITE_SOFT
    cbg.line.fill.background()
    add_tb(sl, x + Inches(0.3), Inches(3.3), Inches(2.9), Inches(0.5), title,
           Pt(22), INK, True, 'Microsoft YaHei')
    add_tb(sl, x + Inches(0.3), Inches(4.0), Inches(2.9), Inches(1.2), desc,
           Pt(13), BODY_LIGHT, False, 'Microsoft YaHei', line_sp=1.6)

foot(sl, "Page 07 · 三大核心职责")

# ============================================================
# SLIDE 8: Work Details (dark)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, INK)
chrome(sl, "Details · 工作详解", page="08 / 14", dark=True)

kicker(sl, "Depth · 深入看", Inches(1), Inches(1.5), Inches(5.5), True)
add_tb(sl, Inches(1), Inches(2.1), Inches(5.5), Inches(0.9), "党员挂牌\n+ 党员查寝",
       Pt(30), PAPER, True, 'Microsoft YaHei')
callout_box(sl, Inches(1), Inches(4.0), Inches(5.5), Inches(1.2),
            "「挂牌的关键不是挂，\n是前期的统筹和收尾的闭环。」",
            "— 我的工作经验", True)

mono(sl, "党员挂牌 · 每学期一次", Inches(7.2), Inches(2.1), Inches(5), True)
body(sl, "提前完成人员分配，挂牌当天分发牌子、床铺贴和工作守则，全程管理纪律，结束后逐一收取照片，确保归档完整。",
     Inches(7.2), Inches(2.6), Inches(5), Inches(0.9), True)
mono(sl, "党员查寝 · 每月一次", Inches(7.2), Inches(4.1), Inches(5), True)
body(sl, "流程与例行查寝相似但更正式：需提前收取按支部分组的签到表，查寝人员每次轮换，主题每次不同，查寝表含心得栏由我们统一收取。",
     Inches(7.2), Inches(4.6), Inches(5), Inches(0.9), True)

foot(sl, "Page 08 · 工作详解", dark=True)

# ============================================================
# SLIDE 9: Routine Inspection Pipeline (light)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, PAPER)
chrome(sl, "Workflow · 工作流", page="09 / 14")

kicker(sl, "Pipeline · 完整流程", Inches(1), Inches(1.3), Inches(11))
h1(sl, "例行查寝，从通知到归档。", Inches(1), Inches(1.8), Inches(11))

mono(sl, "准备阶段 · Preparation", Inches(1), Inches(2.9), Inches(5))
prep = [
    ("01", "发送通知", "周二大课间发布查寝通知与请假表格"),
    ("02", "打印材料", "周三晚前打印签到表、查寝表及评分细则"),
    ("03", "现场布置", "提前摆放桌子、笔、牌子，等待人员到齐"),
    ("04", "点名分组", "管纪律、点名、宣读细则、完成分组分配"),
]
for i, (nb, t, d) in enumerate(prep):
    x = Inches(1 + i * 3.0)
    line(sl, x, Inches(3.3), Inches(2.5))
    mono(sl, nb, x, Inches(3.4), Inches(2.5))
    add_tb(sl, x, Inches(3.7), Inches(2.5), Inches(0.25), t,
           Pt(15), INK, True, 'Microsoft YaHei')
    add_tb(sl, x, Inches(4.05), Inches(2.5), Inches(0.5), d,
           Pt(11), MUTED_DARK, False, 'Microsoft YaHei')

mono(sl, "收尾阶段 · Wrap-up", Inches(1), Inches(5.0), Inches(5))
wrap = [
    ("05", "整理物品", "查寝结束后统一回收物料"),
    ("06", "对接安全部", "查寝表交安全部等出问题表"),
    ("07", "筛选反馈", "优差寝照片发宣传部，做群相册"),
    ("08", "统一归档", "问题表与查寝表装订便于查阅"),
]
for i, (nb, t, d) in enumerate(wrap):
    x = Inches(1 + i * 3.0)
    line(sl, x, Inches(5.4), Inches(2.5))
    mono(sl, nb, x, Inches(5.5), Inches(2.5))
    add_tb(sl, x, Inches(5.8), Inches(2.5), Inches(0.25), t,
           Pt(15), INK, True, 'Microsoft YaHei')
    add_tb(sl, x, Inches(6.15), Inches(2.5), Inches(0.4), d,
           Pt(11), MUTED_DARK, False, 'Microsoft YaHei')

foot(sl, "Page 09 · 查寝全流程")

# ============================================================
# SLIDE 10: Growth (dark)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, INK)
chrome(sl, "Growth · 成长收获", page="10 / 14", dark=True)

kicker(sl, "这一年 · The Year", Inches(1), Inches(1.5), Inches(5.5), True)
add_tb(sl, Inches(1), Inches(2.1), Inches(5.5), Inches(0.9), "在重复里\n长出判断力。",
       Pt(32), PAPER, True, 'Microsoft YaHei')
lead(sl, '一年前我走进党站，只知道「按流程走」。\n一年后我知道：流程是底线，\n但真正重要的是遇到突发情况时的判断。',
     Inches(1), Inches(3.2), Inches(5.5), True)
callout_box(sl, Inches(1), Inches(5.2), Inches(5.5), Inches(0.9),
            "「每次面对的人、宿舍、问题都不一样。\n这就是经验积累的过程。」",
            "— 安全服务部的日常", True)

learned = [
    ("学会了", "沟通", "跟安全部对接、跟宣传部确认、跟组织部协调——跨部门协作是日常。"),
    ("学会了", "统筹", "从人员分配到物料准备，从签到到归档——闭环思维每一次都在强化。"),
    ("学会了", "耐心", "查寝不是走过场，在重复中保持标准不降、不敷衍、不糊弄。"),
]
for i, (lbl, val, note) in enumerate(learned):
    y = Inches(1.8 + i * 1.6)
    mono(sl, lbl, Inches(7.2), y, Inches(2), True)
    add_tb(sl, Inches(7.2), y + Inches(0.2), Inches(3), Inches(0.45), val,
           Pt(24), PAPER, True, 'Microsoft YaHei')
    body(sl, note, Inches(7.2), y + Inches(0.75), Inches(5), Inches(0.5), True)

foot(sl, "Page 10 · 成长收获", dark=True)

# ============================================================
# SLIDE 11: Chapter - Future (hero light)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, PAPER_TINT)
chrome(sl, "第三幕 · 未来展望", page="11 / 14")

kicker(sl, "Chapter III", Inches(1), Inches(2.2), Inches(5))
add_tb(sl, Inches(1), Inches(2.9), Inches(10), Inches(1.1), "下一步",
       Pt(64), INK, True, 'Microsoft YaHei', line_sp=1.0)
lead(sl, "留下来，不是为了重复，而是为了优化。",
     Inches(1), Inches(4.3), Inches(7))
foot(sl, "第三幕 · 未来展望")

# ============================================================
# SLIDE 12: Four Plans (light)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, PAPER)
chrome(sl, "Plan · 发展规划", page="12 / 14")

kicker(sl, "What's Next · 四项承诺", Inches(1), Inches(1.5), Inches(11))
h1(sl, "如果留下，我会做好这四件事。", Inches(1), Inches(2.0), Inches(11))

plans = [
    ("01 · Standard", "标准不降", "每一次查寝、每一次挂牌，质量不打折扣。把细致变成习惯，把负责刻进日常。"),
    ("02 · Growth", "自我提升", "提升沟通能力与应变能力，在查寝实践中积累管理经验。"),
    ("03 · Connect", "联动协作", "主动加强与其他部门的联系——安全部、宣传部、组织部，衔接更紧密。"),
    ("04 · Mentor", "带新传承", '带好新部员，让他们不只是「会做」，更是「理解为什么这么做」。'),
]
for i, (lbl, val, note) in enumerate(plans):
    col = i % 2
    row = i // 2
    x = Inches(1 + col * 5.8)
    y = Inches(3.1 + row * 1.9)
    card(sl, x, y, Inches(5.2), lbl, val, note)

foot(sl, "Page 12 · 未来展望")

# ============================================================
# SLIDE 13: Reflection (dark)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, INK)
chrome(sl, "Reflection · 心得", page="13 / 14", dark=True)

kicker(sl, "写在最后", Inches(1), Inches(1.5), Inches(5), True)
add_tb(sl, Inches(1), Inches(2.2), Inches(10), Inches(2.0),
       '在安全服务部的这一年，每一个工作我都亲手做过，也犯过错、也踩过坑——但正是这些大大小小的错误，让我真正理解了「负责」两个字的重量。',
       Pt(24), PAPER, True, 'Microsoft YaHei', line_sp=1.45)

lead(sl, '感谢学长学姐的包容和指导。感谢部员伙伴的陪伴——\n从陌生到默契，从「按流程做」到「想清楚了再做」。',
     Inches(1), Inches(4.4), Inches(8), True)
mono(sl, "如果我能够继续留在安全服务部，我会带着这份经验，让部门变得更好。",
     Inches(1), Inches(5.8), Inches(8), True)

foot(sl, "Page 13 · 心得", dark=True)

# ============================================================
# SLIDE 14: Thank You (hero dark)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, INK)
chrome(sl, "Fin · 终章", page="14 / 14", dark=True)

kicker(sl, "Thank You", Inches(1), Inches(2.0), Inches(5), True)
add_tb(sl, Inches(1), Inches(2.7), Inches(10), Inches(1.4), "谢谢",
       Pt(72), PAPER, True, 'Microsoft YaHei')
lead(sl, "述职完毕，请各位批评指正。",
     Inches(1), Inches(4.4), Inches(6), True)
mono(sl, "朱轶涵 · 安全服务部 · 25级软件工程二班",
     Inches(1), Inches(5.3), Inches(10), True)

foot(sl, "述职报告 · 安全服务部", "— 朱轶涵 —", True)

# ============================================================
# SAVE
# ============================================================
output_path = r"c:\Users\Mechrevo\Desktop\新建文件夹\述职报告-朱轶涵.pptx"
prs.save(output_path)
print(f"Done! Saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
